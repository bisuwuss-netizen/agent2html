"""
RAG 素材库 - 智能图片匹配
"""
import os
import json
from typing import List, Dict, Optional
from pathlib import Path

# 可选导入 chromadb
try:
    import chromadb
    from chromadb.config import Settings
    CHROMADB_AVAILABLE = True
except ImportError:
    CHROMADB_AVAILABLE = False
    print("⚠️  chromadb 未安装，将使用简化版素材库")


class MaterialLibrary:
    """
    素材库管理 - 基于 ChromaDB 的向量检索
    如果 chromadb 不可用，使用简化的 JSON 存储
    """

    def __init__(self, persist_directory: str = "./data/material_db"):
        """
        初始化素材库

        Args:
            persist_directory: 数据持久化目录
        """
        self.persist_directory = persist_directory
        os.makedirs(persist_directory, exist_ok=True)
        self.use_chromadb = CHROMADB_AVAILABLE
        
        if self.use_chromadb:
            # 使用 ChromaDB
            self.client = chromadb.PersistentClient(
                path=persist_directory,
                settings=Settings(anonymized_telemetry=False)
            )
            self.collection = self.client.get_or_create_collection(
                name="teaching_materials",
                metadata={"description": "高职教育教学素材库"}
            )
            self.materials = None
        else:
            # 使用简化的 JSON 存储
            self.json_path = os.path.join(persist_directory, "materials.json")
            self.materials = self._load_json()
            self.client = None
            self.collection = None
    
    def _load_json(self) -> List[Dict]:
        """加载 JSON 素材库"""
        if os.path.exists(self.json_path):
            with open(self.json_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return []
    
    def _save_json(self):
        """保存 JSON 素材库"""
        with open(self.json_path, 'w', encoding='utf-8') as f:
            json.dump(self.materials, f, ensure_ascii=False, indent=2)

    def add_material(
        self,
        file_path: str,
        description: str,
        tags: List[str],
        metadata: Optional[Dict] = None
    ) -> str:
        """
        添加素材到库中

        Args:
            file_path: 素材文件路径
            description: 素材描述
            tags: 标签列表 (如: ["机械", "车床", "结构图"])
            metadata: 额外元数据 (如: 格式、分辨率、来源等)

        Returns:
            素材ID
        """
        material_id = f"mat_{len(self.collection.get()['ids']) + 1:05d}"

        # 检查文件是否存在
        if not os.path.exists(file_path):
            print(f"⚠️  文件不存在: {file_path}")
            return None

        # 构建元数据
        mat_metadata = {
            "file_path": file_path,
            "tags": ",".join(tags),
            "type": self._get_file_type(file_path),
            **(metadata or {})
        }

        # 添加到 ChromaDB（使用描述作为文档内容）
        self.collection.add(
            ids=[material_id],
            documents=[description],
            metadatas=[mat_metadata]
        )

        print(f"✅ 素材已添加: {material_id} - {description}")
        return material_id

    def search_materials(
        self,
        query: str,
        n_results: int = 5,
        tags_filter: Optional[List[str]] = None
    ) -> List[Dict]:
        """
        搜索相关素材

        Args:
            query: 搜索查询（如："车床主轴结构示意图"）
            n_results: 返回结果数量
            tags_filter: 标签过滤（如: ["机械"]）

        Returns:
            [
                {
                    "id": "mat_00001",
                    "file_path": "/path/to/image.png",
                    "description": "车床主轴结构示意图",
                    "tags": ["机械", "车床"],
                    "score": 0.85
                }
            ]
        """
        # 构建过滤条件
        where_filter = None
        if tags_filter:
            # 注意：ChromaDB 的 where 过滤需要精确匹配
            # 这里简化处理，实际使用时可能需要更复杂的逻辑
            pass

        # 执行搜索
        results = self.collection.query(
            query_texts=[query],
            n_results=n_results,
            where=where_filter
        )

        # 格式化结果
        materials = []
        if results['ids'] and results['ids'][0]:
            for i, material_id in enumerate(results['ids'][0]):
                materials.append({
                    "id": material_id,
                    "file_path": results['metadatas'][0][i].get('file_path'),
                    "description": results['documents'][0][i],
                    "tags": results['metadatas'][0][i].get('tags', '').split(','),
                    "type": results['metadatas'][0][i].get('type'),
                    "score": 1 - results['distances'][0][i] if results.get('distances') else 0
                })

        return materials

    def batch_import_from_ppt(self, ppt_path: str, major: str) -> int:
        """
        从 PPT 批量导入图片素材

        Args:
            ppt_path: PPT 文件路径
            major: 专业类别（用于打标签）

        Returns:
            导入的素材数量
        """
        try:
            from pptx import Presentation

            prs = Presentation(ppt_path)
            count = 0

            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        # 提取图片
                        image = shape.image
                        image_bytes = image.blob

                        # 保存图片
                        img_filename = f"{major}_slide{slide_idx + 1}_{count + 1}.{image.ext}"
                        img_path = os.path.join(self.persist_directory, "images", img_filename)
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)

                        with open(img_path, 'wb') as f:
                            f.write(image_bytes)

                        # 提取描述（从周围文本推测）
                        description = self._extract_image_context(slide, shape)

                        # 添加到素材库
                        self.add_material(
                            file_path=img_path,
                            description=description,
                            tags=[major, f"第{slide_idx + 1}页"],
                            metadata={
                                "source": ppt_path,
                                "slide_number": slide_idx + 1,
                                "format": image.ext
                            }
                        )

                        count += 1

            print(f"✅ 从 {ppt_path} 导入了 {count} 张图片")
            return count

        except Exception as e:
            print(f"❌ PPT 导入失败: {e}")
            return 0

    def _extract_image_context(self, slide, image_shape) -> str:
        """
        从幻灯片中提取图片周围的文字作为描述
        """
        context_texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                context_texts.append(shape.text.strip())

        # 取前100字作为描述
        context = " ".join(context_texts)[:100]
        return context or "教学图片"

    def _get_file_type(self, file_path: str) -> str:
        """
        获取文件类型
        """
        ext = Path(file_path).suffix.lower()
        if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            return "image"
        elif ext in ['.mp4', '.avi', '.mov', '.wmv']:
            return "video"
        elif ext in ['.mp3', '.wav', '.m4a']:
            return "audio"
        else:
            return "other"

    def get_statistics(self) -> Dict:
        """
        获取素材库统计信息
        """
        all_materials = self.collection.get()
        total = len(all_materials['ids'])

        # 按类型统计
        type_count = {}
        tag_count = {}

        for metadata in all_materials['metadatas']:
            mat_type = metadata.get('type', 'unknown')
            type_count[mat_type] = type_count.get(mat_type, 0) + 1

            tags = metadata.get('tags', '').split(',')
            for tag in tags:
                if tag.strip():
                    tag_count[tag.strip()] = tag_count.get(tag.strip(), 0) + 1

        return {
            "total_materials": total,
            "by_type": type_count,
            "by_tag": tag_count
        }

    def clear_all(self):
        """
        清空素材库（谨慎使用）
        """
        self.client.delete_collection("teaching_materials")
        self.collection = self.client.create_collection("teaching_materials")
        print("⚠️  素材库已清空")


# 示例用法
if __name__ == "__main__":
    # 初始化素材库
    library = MaterialLibrary()

    # 添加素材
    library.add_material(
        file_path="./materials/lathe_structure.png",
        description="车床主轴结构示意图，标注了主要部件",
        tags=["机械", "车床", "结构图"],
        metadata={
            "format": "PNG",
            "resolution": "1920x1080",
            "source": "教材配套PPT"
        }
    )

    # 搜索素材
    results = library.search_materials(
        query="车床主轴",
        n_results=3
    )

    print("\n搜索结果:")
    for mat in results:
        print(f"  - {mat['id']}: {mat['description']} (评分: {mat['score']:.2f})")

    # 查看统计
    stats = library.get_statistics()
    print(f"\n素材库统计: {stats}")
